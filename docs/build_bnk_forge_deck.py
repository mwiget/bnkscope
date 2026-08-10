#!/usr/bin/env python3
"""
Build BNK Forge v2 — High-Level Marketing Deck.
F5 template branding (logo, master, theme) + dark-slide aesthetic.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
TEMPLATE = os.path.join(os.path.expanduser("~"),
                        "Code/ai-agent-setup/docs/S3-Overview.pptx")

# ═══════════════════════════════════════════════════════════════════════════════
# F5 palette — dark-slide aesthetic
# ═══════════════════════════════════════════════════════════════════════════════
BG_DARK   = RGBColor(0x0A, 0x0A, 0x1A)
BG_SLIDE  = RGBColor(0x0F, 0x14, 0x2A)

F5_BLUE   = RGBColor(0x2E, 0x6E, 0xE6)
F5_PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
F5_RED    = RGBColor(0xE4, 0x00, 0x2B)
F5_ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
F5_GREEN  = RGBColor(0x10, 0xB9, 0x81)
F5_CYAN   = RGBColor(0x06, 0xB6, 0xD4)

WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GY  = RGBColor(0xBB, 0xBB, 0xCC)
MID_GY    = RGBColor(0x88, 0x88, 0x99)

CARD_BG   = RGBColor(0x16, 0x1E, 0x38)
CARD_DARK = RGBColor(0x0D, 0x11, 0x22)
TBL_HDR   = RGBColor(0x1E, 0x29, 0x4A)
TBL_R1    = RGBColor(0x1A, 0x23, 0x40)
TBL_R2    = RGBColor(0x14, 0x1C, 0x35)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ═══════════════════════════════════════════════════════════════════════════════
# Load F5 template, strip content slides
# ═══════════════════════════════════════════════════════════════════════════════
prs = Presentation(TEMPLATE)

sldIdLst = prs.slides._sldIdLst
for sldId in list(sldIdLst):
    rId = sldId.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    if rId is None:
        rId = sldId.get('r:id')
    if rId:
        try:
            prs.part.drop_rel(rId)
        except KeyError:
            pass
    sldIdLst.remove(sldId)

LAYOUTS = {layout.name: layout for layout in prs.slide_layouts}

def get_layout(name):
    if name in LAYOUTS:
        return LAYOUTS[name]
    for k, v in LAYOUTS.items():
        if name.lower() in k.lower():
            return v
    return LAYOUTS.get("Blank Layout Black", LAYOUTS.get("Blank", list(LAYOUTS.values())[0]))


# ═══════════════════════════════════════════════════════════════════════════════
# Helpers
# ═══════════════════════════════════════════════════════════════════════════════
def _bg(slide, c):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = c

def _rect(s, l, t, w, h, fc, bc=None):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fc
    if bc: sh.line.color.rgb = bc; sh.line.width = Pt(1)
    else:  sh.line.fill.background()
    return sh

def _rrect(s, l, t, w, h, fc, bc=None):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fc
    if bc: sh.line.color.rgb = bc; sh.line.width = Pt(1.5)
    else:  sh.line.fill.background()
    return sh

def _tb(s, l, t, w, h):
    return s.shapes.add_textbox(l, t, w, h)

def _txt(tf, text, sz=18, c=WHITE, b=False, al=PP_ALIGN.LEFT, fn="Calibri"):
    tf.clear(); tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text; p.font.size = Pt(sz); p.font.color.rgb = c
    p.font.bold = b; p.font.name = fn; p.alignment = al
    return p

def _p(tf, text, sz=18, c=WHITE, b=False, al=PP_ALIGN.LEFT, fn="Calibri",
       sb=Pt(4), sa=Pt(2)):
    p = tf.add_paragraph()
    p.text = text; p.font.size = Pt(sz); p.font.color.rgb = c
    p.font.bold = b; p.font.name = fn; p.alignment = al
    p.space_before = sb; p.space_after = sa
    return p

def _bul(tf, text, sz=14, c=WHITE, b=False, fn="Calibri", lvl=0):
    p = tf.add_paragraph()
    p.text = text; p.font.size = Pt(sz); p.font.color.rgb = c
    p.font.bold = b; p.font.name = fn; p.level = lvl
    p.space_before = Pt(3); p.space_after = Pt(1)
    return p

def _table(s, rows, cols, l, t, w, h):
    return s.shapes.add_table(rows, cols, l, t, w, h).table

def _tcell(cell, text, sz=11, c=WHITE, b=False, bg=None, al=PP_ALIGN.LEFT):
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.text = text; p.font.size = Pt(sz); p.font.color.rgb = c
    p.font.bold = b; p.font.name = "Calibri"; p.alignment = al
    cell.text_frame.word_wrap = True
    if bg: cell.fill.solid(); cell.fill.fore_color.rgb = bg
    cell.margin_left = Pt(6); cell.margin_right = Pt(6)
    cell.margin_top = Pt(4); cell.margin_bottom = Pt(4)

def _arrow(s, x, y, w=Inches(0.5), h=Inches(0.22), col=F5_CYAN):
    ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h)
    ar.fill.solid(); ar.fill.fore_color.rgb = col; ar.line.fill.background()
    return ar

def _bar(s):
    _rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.04), F5_BLUE)

def _slide_title(s, title, subtitle=""):
    tb = _tb(s, Inches(0.8), Inches(0.25), Inches(10), Inches(0.6))
    _txt(tb.text_frame, title, sz=30, b=True)
    if subtitle:
        tb = _tb(s, Inches(0.8), Inches(0.8), Inches(10), Inches(0.4))
        _txt(tb.text_frame, subtitle, sz=15, c=F5_CYAN)

def new_slide():
    layout = get_layout("Blank Layout Black")
    s = prs.slides.add_slide(layout)
    _bg(s, BG_SLIDE)
    _bar(s)
    return s

def note(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 0 — Title
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(get_layout("Title Slide: Black Solid"))
for ph in list(slide.placeholders):
    sp = ph._element; sp.getparent().remove(sp)
_bg(slide, BG_DARK)
_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), F5_RED)

tb = _tb(slide, Inches(1.2), Inches(1.2), Inches(10), Inches(1.2))
tf = tb.text_frame
_txt(tf, "BNK Forge", sz=52, b=True)
_p(tf, "The Management Plane for F5 BNK", sz=36, b=True, c=F5_CYAN)

tb = _tb(slide, Inches(1.2), Inches(3.0), Inches(10), Inches(0.5))
_txt(tb.text_frame,
     "Deploy, Operate, Monitor & Evolve — in Minutes, Not Days",
     sz=22, c=LIGHT_GY)

_rect(slide, Inches(1.2), Inches(3.8), Inches(3), Inches(0.04), F5_RED)

tb = _tb(slide, Inches(1.2), Inches(4.15), Inches(10), Inches(1.2))
tf = tb.text_frame
_txt(tf, "From CLI complexity to GUI simplicity", sz=18, c=MID_GY)
_p(tf, "v2.11  |  200+ API Endpoints  |  91 AI-Operable Tools", sz=14, c=MID_GY)

_rect(slide, Inches(0), Inches(7.44), SLIDE_W, Inches(0.06), F5_PURPLE)
note(slide, "Title slide. BNK Forge is the management plane for F5 BIG-IP Next "
     "for Kubernetes. Replaces hours of CLI work with a unified GUI + AI interface.")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — What is BNK Forge?
# ═══════════════════════════════════════════════════════════════════════════════
slide = new_slide()
_slide_title(slide, "What is BNK Forge?",
             "One platform for the entire BNK lifecycle")

# Lifecycle cards — horizontal flow
stages = [
    ("DEPLOY",   "Blueprints, modules,\nOpenTofu execution",       F5_GREEN),
    ("OPERATE",  "Topology, policy,\nHelm, config promotion",     F5_CYAN),
    ("MONITOR",  "Fleet health, metrics,\nconnectivity checks",   F5_BLUE),
    ("DIAGNOSE", "QKView, runbooks,\nTMM debug, licensing",       F5_PURPLE),
]
for i, (title, desc, col) in enumerate(stages):
    x = Inches(0.8 + i * 3.1)
    _rrect(slide, x, Inches(1.4), Inches(2.8), Inches(1.4), CARD_BG, col)
    tb = _tb(slide, x + Inches(0.15), Inches(1.5), Inches(2.5), Inches(1.2))
    tf = tb.text_frame
    _txt(tf, title, sz=16, c=col, b=True, al=PP_ALIGN.CENTER)
    _p(tf, desc, sz=11, c=LIGHT_GY, al=PP_ALIGN.CENTER)

# Arrows
for ax in [Inches(3.75), Inches(6.85), Inches(9.95)]:
    _arrow(slide, ax, Inches(1.9), w=Inches(0.4), h=Inches(0.18))

# Left: Key capabilities
tb = _tb(slide, Inches(0.8), Inches(3.2), Inches(6), Inches(0.3))
_txt(tb.text_frame, "Key Capabilities", sz=16, c=WHITE, b=True)

_rrect(slide, Inches(0.8), Inches(3.55), Inches(5.8), Inches(3.2), CARD_BG, F5_CYAN)
tb = _tb(slide, Inches(1.0), Inches(3.65), Inches(5.4), Inches(3.0))
tf = tb.text_frame
_txt(tf, "", sz=2, c=WHITE)
_bul(tf, "Multi-cluster fleet management from kubeconfigs", sz=12, c=LIGHT_GY)
_bul(tf, "Gateway topology — the central UX for BNK traffic & security", sz=12, c=LIGHT_GY)
_bul(tf, "Infrastructure-as-Code with OpenTofu modules & blueprints", sz=12, c=LIGHT_GY)
_bul(tf, "Full Helm lifecycle — install, upgrade, rollback, uninstall", sz=12, c=LIGHT_GY)
_bul(tf, "Day-2 diagnostics — QKView, TMM debug, runbooks, licensing", sz=12, c=LIGHT_GY)
_bul(tf, "Configuration diff, export, import & promotion across clusters", sz=12, c=LIGHT_GY)
_bul(tf, "DPF integration — 18 CRDs, full infrastructure visibility", sz=12, c=LIGHT_GY)
_bul(tf, "RBAC + multi-user ownership + full audit trail", sz=12, c=LIGHT_GY)

# Right: Architecture summary
tb = _tb(slide, Inches(7.0), Inches(3.2), Inches(6), Inches(0.3))
_txt(tb.text_frame, "Architecture", sz=16, c=WHITE, b=True)

arch_boxes = [
    ("React 18 + Vite", "Modern SPA\nTailwind + shadcn/ui", F5_GREEN),
    ("FastAPI + Celery", "200+ endpoints\nSQLAlchemy + Redis", F5_CYAN),
    ("MCP Server", "91 AI-operable tools\nStreamable HTTP", F5_PURPLE),
]
for i, (title, desc, col) in enumerate(arch_boxes):
    y = Inches(3.6 + i * 1.05)
    _rrect(slide, Inches(7.0), y, Inches(5.5), Inches(0.9), CARD_BG, col)
    tb = _tb(slide, Inches(7.2), y + Inches(0.05), Inches(2.4), Inches(0.8))
    _txt(tb.text_frame, title, sz=13, c=col, b=True)
    tb = _tb(slide, Inches(9.6), y + Inches(0.05), Inches(2.7), Inches(0.8))
    _txt(tb.text_frame, desc, sz=10, c=LIGHT_GY)

note(slide, "BNK Forge covers the full lifecycle: deploy, operate, monitor, diagnose. "
     "Three-tier architecture: React frontend, FastAPI backend, MCP AI server.")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Before vs After (The Problem & The Solution)
# ═══════════════════════════════════════════════════════════════════════════════
slide = new_slide()
_slide_title(slide, "Before vs After BNK Forge",
             "Hours of CLI complexity  →  minutes in a GUI")

# Comparison table — the core message
tbl = _table(slide, 8, 4, Inches(0.8), Inches(1.3), Inches(11.7), Inches(3.2))
for i, h in enumerate(["Task", "Without BNK Forge", "With BNK Forge", "Savings"]):
    _tcell(tbl.cell(0, i), h, sz=11, b=True, c=F5_CYAN, bg=TBL_HDR)

rows = [
    ["Deploy BNK to a cluster",
     "Manual kubectl, Helm,\nTofu CLI — hours of YAML",
     "Blueprint → click Deploy\n→ watch progress",
     "Hours → 5 min"],
    ["Check fleet health",
     "SSH into each cluster,\nrun kubectl per node",
     "Fleet dashboard —\nall clusters at a glance",
     "30 min → 10 sec"],
    ["Troubleshoot BNK issue",
     "Dig through logs, exec\ninto pods, read docs",
     "QKView + TMM debug +\nrunbooks in one UI",
     "Hours → minutes"],
    ["Promote config across\nclusters",
     "Manual export, diff,\ncopy-paste YAML files",
     "Config diff → review →\none-click promote",
     "Days → minutes"],
    ["Manage Helm releases",
     "helm CLI per cluster,\ntrack versions manually",
     "Visual Helm manager\nwith rollback support",
     "Per-cluster → unified"],
    ["Understand traffic flow",
     "Read Gateway CRDs,\npiece together manually",
     "Gateway Topology —\nvisual + interactive",
     "Guesswork → clarity"],
    ["AI-assisted operations",
     "Not possible",
     "91 MCP tools — AI can\nquery, diagnose, operate",
     "N/A → game changer"],
]

for r, row in enumerate(rows):
    bg = TBL_R1 if r % 2 == 0 else TBL_R2
    _tcell(tbl.cell(r+1, 0), row[0], sz=10, c=WHITE,    bg=bg, b=True)
    _tcell(tbl.cell(r+1, 1), row[1], sz=10, c=F5_RED,   bg=bg)
    _tcell(tbl.cell(r+1, 2), row[2], sz=10, c=F5_GREEN, bg=bg)
    _tcell(tbl.cell(r+1, 3), row[3], sz=10, c=F5_ORANGE, bg=bg, b=True)

tbl.columns[0].width = Inches(2.2)
tbl.columns[1].width = Inches(3.3)
tbl.columns[2].width = Inches(3.5)
tbl.columns[3].width = Inches(2.7)

# Bottom: key message cards
cards = [
    ("CLI Expertise Required", "Operators need deep kubectl,\nHelm, and Tofu knowledge", F5_RED),
    ("GUI Simplicity", "Point-and-click workflows\nwith guardrails built in", F5_GREEN),
    ("AI-Operable", "91 MCP tools let AI assistants\noperate your infrastructure", F5_PURPLE),
]
for i, (title, desc, col) in enumerate(cards):
    x = Inches(0.8 + i * 4.1)
    _rrect(slide, x, Inches(4.75), Inches(3.7), Inches(1.5), CARD_BG, col)
    tb = _tb(slide, x + Inches(0.2), Inches(4.85), Inches(3.3), Inches(1.3))
    tf = tb.text_frame
    _txt(tf, title, sz=14, c=col, b=True, al=PP_ALIGN.CENTER)
    _p(tf, desc, sz=11, c=LIGHT_GY, al=PP_ALIGN.CENTER)

# Arrows between cards
for ax in [Inches(4.65), Inches(8.75)]:
    _arrow(slide, ax, Inches(5.3))

note(slide, "This is the core value proposition. Every row shows a real workflow "
     "that goes from hours/days of CLI work to minutes in the GUI.")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Who It's For
# ═══════════════════════════════════════════════════════════════════════════════
slide = new_slide()
_slide_title(slide, "Who It's For",
             "Built for every role in the BNK lifecycle")

# ── Persona 1: Platform Engineer ──
_rrect(slide, Inches(0.8), Inches(1.35), Inches(3.7), Inches(5.1), CARD_BG, F5_GREEN)
tb = _tb(slide, Inches(1.0), Inches(1.5), Inches(3.3), Inches(0.4))
_txt(tb.text_frame, "Platform Engineer", sz=20, c=F5_GREEN, b=True, al=PP_ALIGN.CENTER)
tb = _tb(slide, Inches(1.0), Inches(1.95), Inches(3.3), Inches(0.3))
_txt(tb.text_frame, "\"I build and ship infrastructure\"", sz=11, c=MID_GY,
     al=PP_ALIGN.CENTER, fn="Calibri")
_rect(slide, Inches(1.2), Inches(2.35), Inches(2.9), Inches(0.02), F5_GREEN)

tb = _tb(slide, Inches(1.0), Inches(2.55), Inches(3.3), Inches(3.5))
tf = tb.text_frame
_txt(tf, "", sz=2, c=WHITE)
_bul(tf, "Deploy BNK clusters with blueprints", sz=12, c=LIGHT_GY)
_bul(tf, "   — not 200 lines of YAML", sz=10, c=MID_GY, lvl=1)
_bul(tf, "OpenTofu modules with plan/apply UI", sz=12, c=LIGHT_GY)
_bul(tf, "Helm lifecycle — install to rollback", sz=12, c=LIGHT_GY)
_bul(tf, "Config promotion across clusters", sz=12, c=LIGHT_GY)
_bul(tf, "DPF infrastructure visibility", sz=12, c=LIGHT_GY)
_p(tf, "", sz=6, c=WHITE)
_p(tf, "Hours of CLI scripting", sz=12, c=F5_RED, b=True, al=PP_ALIGN.CENTER)
_p(tf, "→  minutes with guided workflows", sz=12, c=F5_GREEN, b=True, al=PP_ALIGN.CENTER)

# ── Persona 2: BNK Operator ──
_rrect(slide, Inches(4.8), Inches(1.35), Inches(3.7), Inches(5.1), CARD_BG, F5_CYAN)
tb = _tb(slide, Inches(5.0), Inches(1.5), Inches(3.3), Inches(0.4))
_txt(tb.text_frame, "BNK Operator", sz=20, c=F5_CYAN, b=True, al=PP_ALIGN.CENTER)
tb = _tb(slide, Inches(5.0), Inches(1.95), Inches(3.3), Inches(0.3))
_txt(tb.text_frame, "\"I keep it running and fix it fast\"", sz=11, c=MID_GY,
     al=PP_ALIGN.CENTER, fn="Calibri")
_rect(slide, Inches(5.2), Inches(2.35), Inches(2.9), Inches(0.02), F5_CYAN)

tb = _tb(slide, Inches(5.0), Inches(2.55), Inches(3.3), Inches(3.5))
tf = tb.text_frame
_txt(tf, "", sz=2, c=WHITE)
_bul(tf, "Fleet health — all clusters at a glance", sz=12, c=LIGHT_GY)
_bul(tf, "Gateway Topology — visual traffic flow", sz=12, c=LIGHT_GY)
_bul(tf, "QKView, TMM debug & runbooks in one UI", sz=12, c=LIGHT_GY)
_bul(tf, "Connectivity checks & diagnostics", sz=12, c=LIGHT_GY)
_bul(tf, "Licensing, drift detection, snapshots", sz=12, c=LIGHT_GY)
_p(tf, "", sz=6, c=WHITE)
_p(tf, "Hours digging through pods & logs", sz=12, c=F5_RED, b=True, al=PP_ALIGN.CENTER)
_p(tf, "→  answers in clicks, not commands", sz=12, c=F5_CYAN, b=True, al=PP_ALIGN.CENTER)

# ── Persona 3: AI Assistant / Automation ──
_rrect(slide, Inches(8.8), Inches(1.35), Inches(3.7), Inches(5.1), CARD_BG, F5_PURPLE)
tb = _tb(slide, Inches(9.0), Inches(1.5), Inches(3.3), Inches(0.4))
_txt(tb.text_frame, "AI Assistant", sz=20, c=F5_PURPLE, b=True, al=PP_ALIGN.CENTER)
tb = _tb(slide, Inches(9.0), Inches(1.95), Inches(3.3), Inches(0.3))
_txt(tb.text_frame, "\"I automate what humans repeat\"", sz=11, c=MID_GY,
     al=PP_ALIGN.CENTER, fn="Calibri")
_rect(slide, Inches(9.2), Inches(2.35), Inches(2.9), Inches(0.02), F5_PURPLE)

tb = _tb(slide, Inches(9.0), Inches(2.55), Inches(3.3), Inches(3.5))
tf = tb.text_frame
_txt(tf, "", sz=2, c=WHITE)
_bul(tf, "91 governed MCP tools", sz=12, c=LIGHT_GY)
_bul(tf, "   — read, diagnose, operate via API", sz=10, c=MID_GY, lvl=1)
_bul(tf, "Risk-classified (read / mutate / destroy)", sz=12, c=LIGHT_GY)
_bul(tf, "Structured error envelopes for recovery", sz=12, c=LIGHT_GY)
_bul(tf, "Full audit trail of every action", sz=12, c=LIGHT_GY)
_bul(tf, "Observable — invocation telemetry built in", sz=12, c=LIGHT_GY)
_p(tf, "", sz=6, c=WHITE)
_p(tf, "Manual repetition & tribal knowledge", sz=12, c=F5_RED, b=True, al=PP_ALIGN.CENTER)
_p(tf, "→  conversational infrastructure ops", sz=12, c=F5_PURPLE, b=True, al=PP_ALIGN.CENTER)

# ── Bottom: unifying callout ──
_rrect(slide, Inches(0.8), Inches(6.65), Inches(11.7), Inches(0.6), TBL_HDR, F5_ORANGE)
tb = _tb(slide, Inches(1.0), Inches(6.7), Inches(11.3), Inches(0.5))
_txt(tb.text_frame,
     "One platform.  Three audiences.  Same source of truth.",
     sz=16, c=F5_ORANGE, b=True, al=PP_ALIGN.CENTER)

note(slide, "Three personas: Platform Engineer (build/deploy), BNK Operator (day-2 ops/diagnostics), "
     "AI Assistant (91 MCP tools for automated operations). All share the same platform and data.")


# ═══════════════════════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════════════════════
out = os.path.join(SCRIPT_DIR, "BNK_Forge_Overview.pptx")
prs.save(out)
print(f"Saved {len(prs.slides)} slides -> {out}")
